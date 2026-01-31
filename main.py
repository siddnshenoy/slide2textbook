from fastapi import FastAPI, UploadFile, File
from pptx import Presentation
import shutil
import os

app = FastAPI()

UPLOAD_DIR = "uploads"

@app.post("/upload")
async def upload_ppt(file: UploadFile = File(...)):
    file_path = os.path.join(UPLOAD_DIR, file.filename)

    # Save file
    with open(file_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    # Load ppt
    prs = Presentation(file_path)

    slides_text = []

    for i, slide in enumerate(prs.slides):
        slide_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_content.append(shape.text)

        print(f"Slide {i+1}: ", slide_content)

        slides_text.append({
            "slide": i + 1,
            "content": slide_content
        })

